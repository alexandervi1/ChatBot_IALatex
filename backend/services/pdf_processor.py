import re
import logging
import os
from typing import List, Dict, Any
from PyPDF2 import PdfReader
import docx
import pptx
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PDFProcessor: # Manteniendo el nombre por compatibilidad de importación
    """
    Servicio para procesar varios tipos de archivo (PDF, DOCX, PPTX, TXT).
    """
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            add_start_index=True,
        )

    def _clean_text(self, text: str) -> str:
        text = re.sub(r'(\w+)-\n(\w+)', r'\1\2', text)
        text = re.sub(r'(?<=[a-zA-Z])\n(?=[a-zA-Z])', ' ', text)
        text = re.sub(r'[ \t]{2,}', ' ', text)
        
        lines = text.split('\n')
        cleaned_lines = []
        for line in lines:
            stripped = line.strip()
            # Keep lines that are either empty (paragraph breaks) or have meaningful content
            if len(stripped) == 0:
                cleaned_lines.append(stripped)
            elif not stripped.isdigit():
                cleaned_lines.append(stripped)
        
        # Re-join and then normalize multiple newlines down to a maximum of two
        text = '\n'.join(cleaned_lines)
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        return text.strip()

    def _process_pdf_content(self, file_path: str) -> List[Document]:
        reader = PdfReader(file_path)
        docs = []
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text and page_text.strip():
                docs.append(Document(
                    page_content=page_text,
                    metadata={"page": i + 1, "source_file": os.path.basename(file_path)}
                ))
        return docs

    def _process_docx_content(self, file_path: str) -> List[Document]:
        doc = docx.Document(file_path)
        full_text = "\n".join([para.text for para in doc.paragraphs])
        return [Document(page_content=full_text, metadata={"source_file": os.path.basename(file_path)})]

    def _process_pptx_content(self, file_path: str) -> List[Document]:
        pres = pptx.Presentation(file_path)
        full_text = "\n".join([shape.text for slide in pres.slides for shape in slide.shapes if hasattr(shape, "text")])
        return [Document(page_content=full_text, metadata={"source_file": os.path.basename(file_path)})]

    def _process_txt_content(self, file_path: str) -> List[Document]:
        with open(file_path, 'r', encoding='utf-8') as f:
            full_text = f.read()
        return [Document(page_content=full_text, metadata={"source_file": os.path.basename(file_path)})]

    def process_file(self, file_path: str) -> List[Dict[str, Any]]:
        logger.info(f"Iniciando procesamiento del archivo: {file_path}")
        _, file_extension = os.path.splitext(file_path.lower())

        try:
            raw_docs = []
            if file_extension == '.pdf':
                raw_docs = self._process_pdf_content(file_path)
            elif file_extension == '.docx':
                raw_docs = self._process_docx_content(file_path)
            elif file_extension == '.pptx':
                raw_docs = self._process_pptx_content(file_path)
            elif file_extension == '.txt':
                raw_docs = self._process_txt_content(file_path)
            else:
                raise ValueError(f"Formato de archivo no soportado: {file_extension}")

            if not raw_docs:
                raise ValueError("No se pudo extraer texto útil del archivo.")

            # Limpiar el texto de cada documento extraído
            for doc in raw_docs:
                doc.page_content = self._clean_text(doc.page_content)

            logger.info(f"Texto extraído y limpiado de {os.path.basename(file_path)}.")

            text_chunks = self.text_splitter.split_documents(raw_docs)
            
            processed_chunks = []
            for i, chunk in enumerate(text_chunks):
                chunk.metadata["chunk_index"] = i
                chunk.metadata["token_count"] = len(chunk.page_content.split())
                processed_chunks.append({
                    "id": f"{os.path.basename(file_path)}-{i}",
                    "content": chunk.page_content,
                    "metadata": chunk.metadata
                })
            
            logger.info(f"Archivo procesado en {len(processed_chunks)} chunks.")
            return processed_chunks

        except FileNotFoundError:
            logger.error(f"Archivo no encontrado: {file_path}")
            raise
        except Exception as e:
            logger.error(f"Error procesando el archivo '{file_path}': {e}")
            raise Exception(f"El archivo '{os.path.basename(file_path)}' podría estar corrupto o no ser soportado.") from e
